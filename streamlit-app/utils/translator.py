# -*- coding: utf-8 -*-
"""
utils/translator.py
HK-AICOS Phase 2.5G — Document Translation & Conversion Engine

Supports:
  - Text translation via AI (EN→TC, TC→EN)
  - DOCX text extraction + translation
  - XLSX text extraction + translation
  - PDF text extraction + translation
  - Excel → PDF conversion
  - DOCX → PDF conversion
  - Image → PowerPoint (BETA — entry point only)

All outputs use embedded NotoSansTC font for cross-platform Chinese rendering.
Never overwrites source files — always returns bytes for download.
"""

import io
import sys
import traceback
from pathlib import Path
from typing import Optional

# ── Font path (same as report_generator) ─────────────────────────────────────
_BASE_DIR  = Path(__file__).parent.parent
_FONT_PATH = _BASE_DIR / "assets" / "fonts" / "NotoSansTC-Regular.ttf"

# ── ReportLab ─────────────────────────────────────────────────────────────────
try:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import mm
    from reportlab.lib import colors
    from reportlab.platypus import (
        SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, HRFlowable
    )
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    _REPORTLAB_OK = True
except ImportError:
    _REPORTLAB_OK = False

# ── Register font once ────────────────────────────────────────────────────────
_FONT_REGISTERED = False

def _ensure_font() -> str:
    """Register NotoSansTC and return font name. Falls back to Helvetica with warning."""
    global _FONT_REGISTERED
    if not _REPORTLAB_OK:
        return "Helvetica"
    if _FONT_REGISTERED:
        return "NotoSansTC"
    if _FONT_PATH.exists():
        try:
            pdfmetrics.registerFont(TTFont("NotoSansTC", str(_FONT_PATH)))
            _FONT_REGISTERED = True
            return "NotoSansTC"
        except Exception as e:
            print(f"[translator] WARNING: font registration failed: {e}", file=sys.stderr)
    else:
        print(f"[translator] WARNING: font not found at {_FONT_PATH}", file=sys.stderr)
    return "Helvetica"


# ── AI translation helper ─────────────────────────────────────────────────────
def translate_text_via_ai(
    text: str,
    direction: str,          # "en_to_tc" | "tc_to_en"
    ai_client=None,
    ai_provider: str = "anthropic",
    model: str = "",
) -> str:
    """
    Translate text using the configured AI provider.
    Returns translated text, or raises RuntimeError if AI is unavailable.
    """
    if not text or not text.strip():
        return ""

    if direction == "en_to_tc":
        instruction = (
            "你是一位專業的工程文件翻譯員，專門處理香港建築及工程行業文件。\n"
            "請將以下英文文字翻譯成繁體中文。\n"
            "規則：\n"
            "1. 保持原文件結構及段落\n"
            "2. 工程術語使用香港慣用繁體中文\n"
            "3. 不改變原意\n"
            "4. 不加入額外解釋\n"
            "5. 只輸出翻譯結果，不輸出任何說明\n\n"
            f"原文：\n{text}"
        )
    elif direction == "tc_to_en":
        instruction = (
            "You are a professional engineering document translator specialising in "
            "Hong Kong construction and engineering documents.\n"
            "Translate the following Traditional Chinese text into English.\n"
            "Rules:\n"
            "1. Preserve original document structure and paragraphs\n"
            "2. Use standard Hong Kong engineering terminology in English\n"
            "3. Do not alter the original meaning\n"
            "4. Do not add extra explanations\n"
            "5. Output only the translation, no commentary\n\n"
            f"Source text:\n{text}"
        )
    else:
        raise ValueError(f"Unknown translation direction: {direction}")

    if ai_client is None:
        raise RuntimeError("AI client not provided. Please configure API key.")

    try:
        if ai_provider == "anthropic":
            _model = model or "claude-3-5-haiku-20241022"
            response = ai_client.messages.create(
                model=_model,
                max_tokens=4096,
                messages=[{"role": "user", "content": instruction}],
            )
            return response.content[0].text.strip()

        elif ai_provider == "openai":
            _model = model or "gpt-4o-mini"
            response = ai_client.chat.completions.create(
                model=_model,
                messages=[{"role": "user", "content": instruction}],
                max_tokens=4096,
            )
            return response.choices[0].message.content.strip()

        else:
            raise RuntimeError(f"Unsupported AI provider: {ai_provider}")

    except Exception as e:
        raise RuntimeError(f"Translation failed: {e}") from e


# ── DOCX extraction ───────────────────────────────────────────────────────────
def extract_text_from_docx(file_bytes: bytes) -> str:
    """Extract all paragraph text from a DOCX file."""
    try:
        from docx import Document
        doc = Document(io.BytesIO(file_bytes))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return "\n".join(paragraphs)
    except Exception as e:
        raise RuntimeError(f"DOCX extraction failed: {e}") from e


def extract_text_from_xlsx(file_bytes: bytes) -> str:
    """Extract all cell text from an XLSX file."""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)
        lines = []
        for sheet in wb.worksheets:
            lines.append(f"[Sheet: {sheet.title}]")
            for row in sheet.iter_rows(values_only=True):
                row_text = "\t".join(str(c) if c is not None else "" for c in row)
                if row_text.strip():
                    lines.append(row_text)
        return "\n".join(lines)
    except Exception as e:
        raise RuntimeError(f"XLSX extraction failed: {e}") from e


def extract_text_from_pdf(file_bytes: bytes) -> str:
    """Extract text from a PDF file using pypdf."""
    try:
        import pypdf
        reader = pypdf.PdfReader(io.BytesIO(file_bytes))
        pages = []
        for i, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            if text.strip():
                pages.append(f"[Page {i+1}]\n{text.strip()}")
        return "\n\n".join(pages)
    except Exception as e:
        raise RuntimeError(f"PDF extraction failed: {e}") from e


# ── Output builders ───────────────────────────────────────────────────────────
def build_translated_pdf(
    translated_text: str,
    source_filename: str,
    direction: str,
    original_text: str = "",
) -> bytes:
    """
    Build a PDF from translated text with embedded NotoSansTC font.
    Returns PDF bytes.
    """
    if not _REPORTLAB_OK:
        raise RuntimeError("reportlab not installed")

    font_name = _ensure_font()
    buf = io.BytesIO()

    doc = SimpleDocTemplate(
        buf,
        pagesize=A4,
        leftMargin=20*mm,
        rightMargin=20*mm,
        topMargin=20*mm,
        bottomMargin=20*mm,
        title=f"Translation — {source_filename}",
        author="HK-AICOS | Buildway Tech (HK) Limited",
    )

    direction_label = "英文 → 繁體中文" if direction == "en_to_tc" else "繁體中文 → 英文"

    styles = {
        "title": ParagraphStyle(
            "title", fontName=font_name, fontSize=16, leading=22,
            textColor=colors.HexColor("#1a3a5c"), spaceAfter=6,
        ),
        "meta": ParagraphStyle(
            "meta", fontName=font_name, fontSize=9, leading=13,
            textColor=colors.HexColor("#666666"), spaceAfter=4,
        ),
        "heading": ParagraphStyle(
            "heading", fontName=font_name, fontSize=11, leading=16,
            textColor=colors.HexColor("#1a3a5c"), spaceBefore=10, spaceAfter=4,
        ),
        "body": ParagraphStyle(
            "body", fontName=font_name, fontSize=10, leading=16,
            textColor=colors.HexColor("#222222"), spaceAfter=4,
        ),
        "disclaimer": ParagraphStyle(
            "disclaimer", fontName=font_name, fontSize=8, leading=12,
            textColor=colors.HexColor("#888888"),
        ),
    }

    story = []

    # Header
    story.append(Paragraph("HK-AICOS 文件翻譯", styles["title"]))
    story.append(Paragraph(f"翻譯方向：{direction_label}", styles["meta"]))
    story.append(Paragraph(f"來源文件：{source_filename}", styles["meta"]))
    story.append(HRFlowable(width="100%", thickness=1, color=colors.HexColor("#c9a84c"), spaceAfter=8))

    # Translated content
    story.append(Paragraph("翻譯結果", styles["heading"]))
    for line in translated_text.split("\n"):
        line = line.strip()
        if line:
            # Escape XML special chars
            safe = line.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
            story.append(Paragraph(safe, styles["body"]))
        else:
            story.append(Spacer(1, 4))

    story.append(Spacer(1, 12))
    story.append(HRFlowable(width="100%", thickness=0.5, color=colors.HexColor("#cccccc"), spaceAfter=6))
    story.append(Paragraph(
        "本翻譯由 HK-AICOS AI 系統生成，僅供參考。工程術語及法律文件請由專業人士確認。"
        " | Buildway Tech (HK) Limited",
        styles["disclaimer"],
    ))

    doc.build(story)
    return buf.getvalue()


def build_translated_docx(
    translated_text: str,
    source_filename: str,
    direction: str,
) -> bytes:
    """Build a DOCX from translated text. Returns DOCX bytes."""
    try:
        from docx import Document
        from docx.shared import Pt, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH
    except ImportError:
        raise RuntimeError("python-docx not installed")

    direction_label = "英文 → 繁體中文" if direction == "en_to_tc" else "繁體中文 → 英文"

    doc = Document()

    # Title
    title_para = doc.add_heading("HK-AICOS 文件翻譯", level=1)
    title_para.runs[0].font.color.rgb = RGBColor(0x1a, 0x3a, 0x5c)

    doc.add_paragraph(f"翻譯方向：{direction_label}")
    doc.add_paragraph(f"來源文件：{source_filename}")
    doc.add_paragraph("─" * 40)

    doc.add_heading("翻譯結果", level=2)

    for line in translated_text.split("\n"):
        doc.add_paragraph(line)

    doc.add_paragraph("─" * 40)
    disclaimer = doc.add_paragraph(
        "本翻譯由 HK-AICOS AI 系統生成，僅供參考。工程術語及法律文件請由專業人士確認。"
        " | Buildway Tech (HK) Limited"
    )
    disclaimer.runs[0].font.size = Pt(8)
    disclaimer.runs[0].font.color.rgb = RGBColor(0x88, 0x88, 0x88)

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def build_translated_txt(translated_text: str, source_filename: str, direction: str) -> bytes:
    """Build a plain UTF-8 TXT from translated text."""
    direction_label = "英文 → 繁體中文" if direction == "en_to_tc" else "繁體中文 → 英文"
    header = (
        f"HK-AICOS 文件翻譯\n"
        f"翻譯方向：{direction_label}\n"
        f"來源文件：{source_filename}\n"
        f"{'─' * 40}\n\n"
    )
    footer = (
        f"\n{'─' * 40}\n"
        "本翻譯由 HK-AICOS AI 系統生成，僅供參考。\n"
        "Buildway Tech (HK) Limited\n"
    )
    return (header + translated_text + footer).encode("utf-8")


# ── Excel → PDF ───────────────────────────────────────────────────────────────
def excel_to_pdf(file_bytes: bytes, source_filename: str) -> bytes:
    """
    Convert XLSX to PDF using reportlab.
    Renders each sheet as a table in the PDF.
    """
    if not _REPORTLAB_OK:
        raise RuntimeError("reportlab not installed")

    try:
        import openpyxl
    except ImportError:
        raise RuntimeError("openpyxl not installed")

    font_name = _ensure_font()
    wb = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)

    buf = io.BytesIO()
    doc = SimpleDocTemplate(
        buf,
        pagesize=A4,
        leftMargin=15*mm,
        rightMargin=15*mm,
        topMargin=15*mm,
        bottomMargin=15*mm,
        title=f"Excel Export — {source_filename}",
        author="HK-AICOS | Buildway Tech (HK) Limited",
    )

    title_style = ParagraphStyle(
        "title", fontName=font_name, fontSize=14, leading=20,
        textColor=colors.HexColor("#1a3a5c"), spaceAfter=4,
    )
    sheet_style = ParagraphStyle(
        "sheet", fontName=font_name, fontSize=11, leading=16,
        textColor=colors.HexColor("#2d5a8e"), spaceBefore=8, spaceAfter=4,
    )
    cell_style = ParagraphStyle(
        "cell", fontName=font_name, fontSize=8, leading=11,
        textColor=colors.HexColor("#222222"),
    )
    disclaimer_style = ParagraphStyle(
        "disclaimer", fontName=font_name, fontSize=8, leading=12,
        textColor=colors.HexColor("#888888"),
    )

    story = []
    story.append(Paragraph(f"Excel 轉換：{source_filename}", title_style))
    story.append(HRFlowable(width="100%", thickness=1, color=colors.HexColor("#c9a84c"), spaceAfter=6))

    for sheet in wb.worksheets:
        story.append(Paragraph(f"工作表：{sheet.title}", sheet_style))

        rows = list(sheet.iter_rows(values_only=True))
        if not rows:
            story.append(Paragraph("（此工作表為空）", cell_style))
            continue

        # Build table data — cap columns to avoid overflow
        MAX_COLS = 10
        table_data = []
        for row in rows:
            row_cells = []
            for cell in row[:MAX_COLS]:
                val = str(cell) if cell is not None else ""
                safe = val.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
                row_cells.append(Paragraph(safe[:100], cell_style))
            # Pad if fewer columns
            while len(row_cells) < min(MAX_COLS, len(rows[0])):
                row_cells.append(Paragraph("", cell_style))
            table_data.append(row_cells)

        if not table_data:
            continue

        num_cols = len(table_data[0])
        available_width = A4[0] - 30*mm
        col_width = available_width / num_cols

        tbl = Table(table_data, colWidths=[col_width] * num_cols, repeatRows=1)
        tbl.setStyle(TableStyle([
            ("BACKGROUND",  (0, 0), (-1, 0),  colors.HexColor("#1a3a5c")),
            ("TEXTCOLOR",   (0, 0), (-1, 0),  colors.white),
            ("FONTNAME",    (0, 0), (-1, -1), font_name),
            ("FONTSIZE",    (0, 0), (-1, -1), 8),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#f5f7fa")]),
            ("GRID",        (0, 0), (-1, -1), 0.3, colors.HexColor("#cccccc")),
            ("VALIGN",      (0, 0), (-1, -1), "TOP"),
            ("TOPPADDING",  (0, 0), (-1, -1), 3),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 3),
        ]))
        story.append(tbl)
        story.append(Spacer(1, 8))

    story.append(HRFlowable(width="100%", thickness=0.5, color=colors.HexColor("#cccccc"), spaceAfter=4))
    story.append(Paragraph(
        f"由 HK-AICOS 轉換 | Buildway Tech (HK) Limited",
        disclaimer_style,
    ))

    doc.build(story)
    return buf.getvalue()


# ── DOCX → PDF ────────────────────────────────────────────────────────────────
def docx_to_pdf(file_bytes: bytes, source_filename: str) -> bytes:
    """
    Convert DOCX to PDF using reportlab.
    Extracts paragraphs and renders them with embedded font.
    """
    if not _REPORTLAB_OK:
        raise RuntimeError("reportlab not installed")

    try:
        from docx import Document
    except ImportError:
        raise RuntimeError("python-docx not installed")

    font_name = _ensure_font()
    doc_in = Document(io.BytesIO(file_bytes))

    buf = io.BytesIO()
    doc_out = SimpleDocTemplate(
        buf,
        pagesize=A4,
        leftMargin=20*mm,
        rightMargin=20*mm,
        topMargin=20*mm,
        bottomMargin=20*mm,
        title=f"DOCX Export — {source_filename}",
        author="HK-AICOS | Buildway Tech (HK) Limited",
    )

    title_style = ParagraphStyle(
        "title", fontName=font_name, fontSize=14, leading=20,
        textColor=colors.HexColor("#1a3a5c"), spaceAfter=4,
    )
    h1_style = ParagraphStyle(
        "h1", fontName=font_name, fontSize=13, leading=18,
        textColor=colors.HexColor("#1a3a5c"), spaceBefore=8, spaceAfter=4,
    )
    h2_style = ParagraphStyle(
        "h2", fontName=font_name, fontSize=11, leading=16,
        textColor=colors.HexColor("#2d5a8e"), spaceBefore=6, spaceAfter=3,
    )
    body_style = ParagraphStyle(
        "body", fontName=font_name, fontSize=10, leading=16,
        textColor=colors.HexColor("#222222"), spaceAfter=3,
    )
    disclaimer_style = ParagraphStyle(
        "disclaimer", fontName=font_name, fontSize=8, leading=12,
        textColor=colors.HexColor("#888888"),
    )

    story = []
    story.append(Paragraph(f"文件轉換：{source_filename}", title_style))
    story.append(HRFlowable(width="100%", thickness=1, color=colors.HexColor("#c9a84c"), spaceAfter=8))

    for para in doc_in.paragraphs:
        text = para.text.strip()
        if not text:
            story.append(Spacer(1, 4))
            continue

        safe = text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
        style_name = para.style.name if para.style else ""

        if "Heading 1" in style_name or style_name.startswith("Title"):
            story.append(Paragraph(safe, h1_style))
        elif "Heading 2" in style_name:
            story.append(Paragraph(safe, h2_style))
        else:
            story.append(Paragraph(safe, body_style))

    story.append(Spacer(1, 12))
    story.append(HRFlowable(width="100%", thickness=0.5, color=colors.HexColor("#cccccc"), spaceAfter=4))
    story.append(Paragraph(
        "由 HK-AICOS 轉換 | Buildway Tech (HK) Limited",
        disclaimer_style,
    ))

    doc_out.build(story)
    return buf.getvalue()


# ── Image → PowerPoint (BETA) ─────────────────────────────────────────────────
def image_to_pptx_beta(file_bytes: bytes, source_filename: str) -> bytes:
    """
    BETA: Convert image to a single-slide PowerPoint.
    Requires python-pptx. Returns PPTX bytes.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor as PptxRGB
    except ImportError:
        raise RuntimeError(
            "python-pptx not installed. "
            "此功能為 Beta 版，需要安裝 python-pptx。"
        )

    prs = Presentation()
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    # Add image
    img_stream = io.BytesIO(file_bytes)
    slide.shapes.add_picture(img_stream, Inches(0.5), Inches(0.5), width=Inches(9))

    # Add caption
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.5))
    tf = txBox.text_frame
    tf.text = f"來源：{source_filename} | HK-AICOS Beta | Buildway Tech (HK) Limited"
    tf.paragraphs[0].runs[0].font.size = Pt(9)
    tf.paragraphs[0].runs[0].font.color.rgb = PptxRGB(0x88, 0x88, 0x88)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ── Translator Agent personality (for prompt injection) ───────────────────────
TRANSLATOR_AGENT_INSTRUCTION = """你是 HK-AICOS 文件翻譯 Agent，專門處理香港建築及工程行業文件翻譯。

性格：準確、保守、工程術語優先、香港工程語境、不亂改原意。

職責：
- 英文轉繁體中文（香港工程術語）
- 繁體中文轉英文（標準工程英文）
- 保持原文件結構及段落
- 工程術語使用香港慣用表達
- 不加入額外解釋或意見
- 不改變原意

輸出規則：
- 只輸出翻譯結果
- 不輸出任何說明、前言或後記
- 保持原文段落結構
- 數字、代號、圖則編號保持原樣
"""
